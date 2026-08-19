#!/usr/bin/env python3
"""Route-level regression coverage for features/routes/ingest.py (viewer-audits-upgrades audit:
7 verified findings). Starts the real server against the deterministic fixture DB and hits the
actual HTTP routes -- same pattern as test_hardening.py / test_routes.py -- so a regression in the
route wiring itself (not just the underlying feature module, which some of these already have unit
coverage for) fails red.

Covers:
  - POST /api/ingest_scan: duplicate-by-filename detection actually works now (documents.path ->
    os.path.basename(), not the nonexistent documents.filename column) + the VIEWER_INGEST_ROOTS
    fence on this route.
  - GET  /api/ingest_preview: gated by _exposed_read_guard() like its /api/ingest_status sibling
    (previously the one route in this family that leaked real host paths/filenames unauthenticated
    in exposed mode).
  - POST /api/ingest: the route that actually kicks off a background crawl subprocess, exercised
    through the real HTTP path for the first time -- the VIEWER_INGEST_ROOTS fence and the
    already-running guard, both via the route (subprocess.Popen + safeguard mocked, same technique
    test_features.py's ingest_start race test already established).
  - POST /api/airgap_manifest: builds a real signed manifest for a real folder + enforces the fence
    (this route used to bypass the fence entirely).
  - POST /api/airgap_verify: accept / tamper-reject / wrong-secret-reject + enforces the fence (this
    route was the one specifically missing it after r_airgap_manifest was fixed).
  - POST /api/form_2404 + /api/form_2407: a filled payload's data actually lands in the rendered
    PDF text (the existing GET/empty-body sweep only proves the blank form doesn't 5xx).

Pure stdlib except PyMuPDF (fitz/pymupdf), imported lazily and skipped if unavailable for the PDF
text-extraction checks -- same convention as test_medium_fixes.py / test_extraction.py.
RUN ON WINDOWS / a coherent env (imports viewer_app)."""
import json
import os
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import threading
import time
import urllib.error
import urllib.parse
import urllib.request

HERE = os.path.dirname(os.path.abspath(__file__))
ENGINE = os.path.dirname(HERE)
sys.path.insert(0, ENGINE); sys.path.insert(0, HERE)
import fixture                                            # noqa: E402

PORT = 8894
BASE = "http://127.0.0.1:%d" % PORT


def _req(path, data=None, hdrs=None):
    body = json.dumps(data).encode() if data is not None else None
    h = dict(hdrs or {})
    if body is not None:
        h.setdefault("Content-Type", "application/json")
    r = urllib.request.Request(BASE + path, data=body, headers=h)
    try:
        with urllib.request.urlopen(r, timeout=15) as x:
            return x.status, x.read()
    except urllib.error.HTTPError as e:
        return e.code, e.read()
    except Exception as e:
        return -1, str(e).encode()


def _json(b):
    try:
        return json.loads(b)
    except Exception:
        return {}


def main():
    tmp = tempfile.mkdtemp(prefix="ingest_routes_")
    db, _corr = fixture.build(tmp)
    import viewer_app as V
    V.DB_PATH = db; V.INDEX_DIR = os.path.dirname(db)
    from http.server import ThreadingHTTPServer
    srv = ThreadingHTTPServer(("127.0.0.1", PORT), V.Handler)
    threading.Thread(target=srv.serve_forever, daemon=True).start(); time.sleep(0.4)

    tests = []

    def check(name, cond):
        tests.append((name, bool(cond)))

    # Deterministic baseline: no ingest-root fence configured unless a section below sets one.
    _orig_roots = os.environ.pop("VIEWER_INGEST_ROOTS", None)

    try:
        # =====================================================================================
        # POST /api/ingest_scan -- duplicate-by-filename detection (the `filename` column fix)
        # =====================================================================================
        scan_dir = tempfile.mkdtemp(prefix="ingest_scan_")
        dup_name = "TM-9-9999-OLD-10.pdf"
        open(os.path.join(scan_dir, dup_name), "wb").write(b"%PDF-1.4 re-downloaded copy" + b"\x00" * 200)
        open(os.path.join(scan_dir, "TM-9-0000-NEW-10.pdf"), "wb").write(b"%PDF-1.4 brand new manual" + b"\x00" * 200)
        con = sqlite3.connect(db)
        # Simulate a manual already indexed under a DIFFERENT path (a prior ingest run, a different
        # drive/mount) but the SAME filename -- the exact "re-download / re-mounted USB / re-copy"
        # scenario the finding describes.
        con.execute("INSERT INTO documents(id, path) VALUES (999, ?)",
                    (os.path.join("Z:\\some\\prior\\ingest\\path", dup_name),))
        con.commit(); con.close()

        c, b = _req("/api/ingest_scan", {"folder": scan_dir})
        r = _json(b)
        check("ingest_scan -> 200", c == 200)
        check("ingest_scan finds both files", r.get("found") == 2)
        check("ingest_scan detects the already-indexed filename as a duplicate (bug fix)",
              r.get("counts", {}).get("duplicate") == 1)
        check("ingest_scan counts only the truly-new file as new",
              r.get("counts", {}).get("new") == 1)
        dup_reported = {d.get("name") for d in r.get("duplicate", [])}
        check("ingest_scan flags the correct file as the duplicate", dup_name in dup_reported)
        new_reported = {d.get("name") for d in r.get("new", [])}
        check("ingest_scan does NOT flag the new file as duplicate", "TM-9-0000-NEW-10.pdf" in new_reported)

        # VIEWER_INGEST_ROOTS fence must still be enforced on this route (finding: zero coverage).
        fence_root = tempfile.mkdtemp(prefix="ingest_scan_fenceroot_")
        outside_dir = tempfile.mkdtemp(prefix="ingest_scan_outside_")
        os.environ["VIEWER_INGEST_ROOTS"] = fence_root
        try:
            c, b = _req("/api/ingest_scan", {"folder": outside_dir})
            check("ingest_scan fence rejects a folder outside VIEWER_INGEST_ROOTS -> 400", c == 400)
            inside_dir = os.path.join(fence_root, "inside")
            os.makedirs(inside_dir, exist_ok=True)
            open(os.path.join(inside_dir, "ok.pdf"), "wb").write(b"%PDF-1.4 x")
            c, b = _req("/api/ingest_scan", {"folder": inside_dir})
            check("ingest_scan fence allows a folder inside VIEWER_INGEST_ROOTS -> 200", c == 200)
        finally:
            os.environ.pop("VIEWER_INGEST_ROOTS", None)

        # =====================================================================================
        # GET /api/ingest_preview -- must be gated by _exposed_read_guard() (security fix)
        # =====================================================================================
        preview_dir = tempfile.mkdtemp(prefix="ingest_preview_exposed_")
        open(os.path.join(preview_dir, "SECRET_MANUAL_ROUTE_TEST.pdf"), "wb").write(b"%PDF-1.4 x")
        old_exposed, old_token = V._EXPOSED, V._AUTH_TOKEN
        try:
            V._EXPOSED = True
            V._AUTH_TOKEN = "route-test-token-2026"
            c, b = _req("/api/ingest_preview?path=" + urllib.parse.quote(preview_dir))
            check("ingest_preview exposed + no token -> 401 (fix)", c == 401)
            check("ingest_preview exposed + no token: real folder path not leaked",
                  preview_dir.encode() not in b)
            check("ingest_preview exposed + no token: real filename not leaked",
                  b"SECRET_MANUAL_ROUTE_TEST" not in b)

            c, b = _req("/api/ingest_preview?path=" + urllib.parse.quote(preview_dir),
                        hdrs={"X-Viewer-Token": "route-test-token-2026"})
            check("ingest_preview exposed + valid token -> 200", c == 200)
            check("ingest_preview exposed + valid token still functions (not over-blocked)",
                  b"SECRET_MANUAL_ROUTE_TEST" in b)

            # sanity: its already-guarded sibling behaves identically (same gate, not a re-implementation)
            c, b = _req("/api/ingest_status")
            check("ingest_status exposed + no token -> 401 (sanity/parity)", c == 401)
        finally:
            V._EXPOSED, V._AUTH_TOKEN = old_exposed, old_token

        # =====================================================================================
        # POST /api/ingest -- the route was never invoked anywhere; exercise the real HTTP path
        # =====================================================================================
        from features import ingest_feature as _ingest_mod
        import subprocess as _subprocess_mod
        _real_popen = _subprocess_mod.Popen
        _real_sg_mod = sys.modules.get("safeguard")
        _popen_calls = []

        class _FakeProc:
            def __init__(self, running):
                self._running = running
            def poll(self):
                return None if self._running else 0

        class _FakeSafeguard:
            def snapshot(self, *a, **kw):
                return ("SNAP_test", {})

        def _fake_popen(cmd, **kw):
            _popen_calls.append(cmd)
            return _FakeProc(running=False)

        ingest_dir = tempfile.mkdtemp(prefix="ingest_post_")
        open(os.path.join(ingest_dir, "some.pdf"), "wb").write(b"%PDF-1.4 x")
        fence_root2 = tempfile.mkdtemp(prefix="ingest_post_fenceroot_")
        try:
            _subprocess_mod.Popen = _fake_popen
            sys.modules["safeguard"] = _FakeSafeguard()
            _ingest_mod._INGEST = {"proc": None, "path": "", "started": 0.0}

            # fence: folder outside the configured roots -> ok:false, no subprocess ever launched.
            # (p_ingest always answers HTTP 200 -- the fence result travels in the JSON body, same
            # as ingest_start()'s own contract -- so we assert on the body, not the status code.)
            os.environ["VIEWER_INGEST_ROOTS"] = fence_root2
            c, b = _req("/api/ingest", {"path": ingest_dir})
            r = _json(b)
            check("POST /api/ingest fence rejects outside folder (ok:false)", c == 200 and r.get("ok") is False)
            check("POST /api/ingest fence: error mentions ingest roots",
                  "ingest root" in (r.get("error") or "").lower())
            check("POST /api/ingest fence: no subprocess launched", len(_popen_calls) == 0)
            os.environ.pop("VIEWER_INGEST_ROOTS", None)

            # success path: folder allowed (no fence configured) -> ok:true, started:true, the FULL
            # in-app pipeline launched (crawl + OCR + parts extraction, one job -- viewer_ingest.py's
            # 'run' subcommand, not the old crawl-only 'crawl' -- this is the actual "properly scans
            # AND OCRs documents in-app" fix: OCR used to be an explicit separate step (run
            # START-OCR-NOW.bat yourself) after this route finished; now it's part of the same job).
            c, b = _req("/api/ingest", {"path": ingest_dir})
            r = _json(b)
            check("POST /api/ingest starts the run -> ok + started", c == 200 and r.get("ok") is True and r.get("started") is True)
            check("POST /api/ingest actually launched a subprocess", len(_popen_calls) == 1)
            check("POST /api/ingest launches viewer_ingest.py's 'run' subcommand (crawl+OCR+parts), not bare 'crawl'",
                  "run" in _popen_calls[0] and "crawl" not in _popen_calls[0])

            # already-running guard, exercised through the actual route (not just ingest_start()
            # called directly, which test_features.py already covers at the module level).
            _ingest_mod._INGEST = {"proc": _FakeProc(running=True), "path": ingest_dir, "started": time.time(), "kind": "scan"}
            c, b = _req("/api/ingest", {"path": ingest_dir})
            r = _json(b)
            check("POST /api/ingest already-running -> ok:false", c == 200 and r.get("ok") is False)
            check("POST /api/ingest already-running error message", "already in progress" in (r.get("error") or ""))
            check("POST /api/ingest already-running: no second subprocess launched", len(_popen_calls) == 1)

            # =================================================================================
            # POST /api/ocr_backlog_start -- the second half of the same in-app job model: finish
            # OCR on whatever's already queued (no folder path needed), sharing the SAME one-
            # job-at-a-time lock as /api/ingest above (both go through ingest_feature._launch()).
            # Requires confirm:true -- unlike /api/ingest, this route has no required parameter
            # that would naturally reject a bare/empty POST before it does anything, so an empty
            # body must be a clean no-op (caught live: test_routes.py's generic empty-body POST
            # sweep was silently launching a REAL subprocess + taking a REAL safeguard snapshot
            # every time this suite ran, before this gate existed).
            # =================================================================================
            _ingest_mod._INGEST = {"proc": None, "path": "", "started": 0.0, "kind": None}
            _popen_calls_before = len(_popen_calls)   # already 1 by here, from the earlier /api/ingest success-path call above
            c, b = _req("/api/ocr_backlog_start", {})
            r = _json(b)
            check("POST /api/ocr_backlog_start bare/empty body -> 400, no launch",
                  c == 400 and len(_popen_calls) == _popen_calls_before)

            c, b = _req("/api/ocr_backlog_start", {"confirm": True})
            r = _json(b)
            check("POST /api/ocr_backlog_start confirm:true -> ok + started", c == 200 and r.get("ok") is True and r.get("started") is True)
            check("POST /api/ocr_backlog_start launches viewer_ingest.py's 'ocrall' subcommand",
                  "ocrall" in _popen_calls[-1])
            check("POST /api/ocr_backlog_start needs no --root (no folder involved)",
                  "--root" not in _popen_calls[-1])

            # shares the lock with /api/ingest: a backlog job already running blocks a NEW /api/ingest
            # call too (and vice versa -- proven above) -- genuinely one job at a time against this DB.
            _ingest_mod._INGEST = {"proc": _FakeProc(running=True), "path": "", "started": time.time(), "kind": "ocr_backlog"}
            c, b = _req("/api/ingest", {"path": ingest_dir})
            r = _json(b)
            check("POST /api/ingest blocked while an OCR-backlog job is running (shared lock)",
                  c == 200 and r.get("ok") is False and "already in progress" in (r.get("error") or ""))

            # =================================================================================
            # POST /api/ingest_upload -- a single file's bytes straight from the browser (drag-
            # and-drop). Validation-only here (Popen/safeguard still mocked) -- the real end-to-
            # end upload-through-the-real-HTTP-route-to-a-real-subprocess check runs unmocked,
            # after this whole try/finally, same split as the progress-stamping e2e section below.
            # =================================================================================
            _ingest_mod._INGEST = {"proc": None, "path": "", "started": 0.0, "kind": None}
            _popen_before = len(_popen_calls)
            import base64 as _b64

            # Discovery Engine: ingest_upload() now accepts images/.txt/.html/.docx/.xlsx/.pptx/.rtf,
            # matching viewer_ingest.py's index_other() -- .doc (the legacy PRE-2007 binary Word
            # format) is the real still-unsupported case now (no good dependency-light reader exists
            # for it; see engine/office.py's own module docstring), not .docx.
            c, b = _req("/api/ingest_upload", {"filename": "manual.doc", "data": _b64.b64encode(b"hi").decode()})
            r = _json(b)
            check("ingest_upload rejects an unsupported extension (.doc, legacy binary)", c == 200 and r.get("ok") is False
                  and "unsupported" in (r.get("error") or "").lower())
            check("ingest_upload unsupported ext: no subprocess launched", len(_popen_calls) == _popen_before)

            c, b = _req("/api/ingest_upload", {"filename": "notes.txt", "data": _b64.b64encode(b"Gasket length 2.0 in.").decode()})
            r = _json(b)
            check("ingest_upload accepts a .txt file -> ok + started",
                  c == 200 and r.get("ok") is True and r.get("started") is True)
            check("ingest_upload .txt launched a subprocess", len(_popen_calls) == _popen_before + 1)
            _popen_before = len(_popen_calls)   # rebase -- the .txt accept above is a real launch, not a rejection

            c, b = _req("/api/ingest_upload", {"filename": "notreally.png", "data": _b64.b64encode(b"not real image bytes").decode()})
            r = _json(b)
            check("ingest_upload rejects an image extension whose bytes don't actually decode as an image",
                  c == 200 and r.get("ok") is False)

            c, b = _req("/api/ingest_upload", {"filename": "", "data": _b64.b64encode(b"%PDF-1.4 x").decode()})
            r = _json(b)
            check("ingest_upload rejects an empty filename", c == 200 and r.get("ok") is False)

            c, b = _req("/api/ingest_upload", {"filename": "bad.pdf", "data": "not-valid-base64!!!"})
            r = _json(b)
            check("ingest_upload rejects undecodable base64", c == 200 and r.get("ok") is False)

            c, b = _req("/api/ingest_upload", {"filename": "notreally.pdf", "data": _b64.b64encode(b"NOT A PDF AT ALL").decode()})
            r = _json(b)
            check("ingest_upload rejects bytes without a %PDF header", c == 200 and r.get("ok") is False
                  and "PDF" in (r.get("error") or ""))
            check("ingest_upload bad-header: no subprocess launched", len(_popen_calls) == _popen_before)

            oversized_stub = b"%PDF-1.4 " + b"\x00" * 10
            orig_cap = _ingest_mod.UPLOAD_MAX_BYTES
            try:
                _ingest_mod.UPLOAD_MAX_BYTES = 5     # force the size-cap branch without a real 150MB payload
                c, b = _req("/api/ingest_upload", {"filename": "big.pdf", "data": _b64.b64encode(oversized_stub).decode()})
                r = _json(b)
                check("ingest_upload rejects a file over the size cap", c == 200 and r.get("ok") is False
                      and "large" in (r.get("error") or "").lower())
            finally:
                _ingest_mod.UPLOAD_MAX_BYTES = orig_cap

            # a genuinely valid upload: real %PDF bytes, saved to disk, subprocess launched with the
            # 'run' subcommand against the uploads/ folder (not the folder-scan path -- no --root
            # value equal to a caller-supplied path; the uploads dir is server-owned).
            good_pdf = b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj<</Type/Catalog>>endobj\n%%EOF"
            c, b = _req("/api/ingest_upload", {"filename": "TM-UPLOAD-TEST.pdf",
                                                "data": "data:application/pdf;base64," + _b64.b64encode(good_pdf).decode()})
            r = _json(b)
            check("ingest_upload accepts a real PDF (incl. data: URI prefix stripped) -> ok + started",
                  c == 200 and r.get("ok") is True and r.get("started") is True)
            check("ingest_upload echoes the saved filename back", r.get("filename") == "TM-UPLOAD-TEST.pdf")
            check("ingest_upload launched a subprocess", len(_popen_calls) == _popen_before + 1)
            check("ingest_upload launches the 'run' subcommand", "run" in _popen_calls[-1])
            uploads_dir = os.path.join(os.path.dirname(db), "uploads")
            saved_path = os.path.join(uploads_dir, "TM-UPLOAD-TEST.pdf")
            check("ingest_upload actually wrote the file to disk", os.path.exists(saved_path))
            if os.path.exists(saved_path):
                with open(saved_path, "rb") as f:
                    check("ingest_upload wrote the exact decoded bytes (no corruption)", f.read() == good_pdf)

            # a second upload of the SAME filename must not silently overwrite the first (R1/R6).
            _ingest_mod._INGEST = {"proc": None, "path": "", "started": 0.0, "kind": None}
            other_pdf = b"%PDF-1.4\n%DIFFERENT CONTENT\n%%EOF"
            c, b = _req("/api/ingest_upload", {"filename": "TM-UPLOAD-TEST.pdf", "data": _b64.b64encode(other_pdf).decode()})
            r = _json(b)
            check("ingest_upload same-name re-upload -> ok, suffixed, original untouched",
                  c == 200 and r.get("ok") is True and r.get("filename") != "TM-UPLOAD-TEST.pdf")
            with open(saved_path, "rb") as f:
                check("ingest_upload: the ORIGINAL file's bytes are unchanged after the re-upload", f.read() == good_pdf)
        finally:
            _subprocess_mod.Popen = _real_popen
            if _real_sg_mod is not None: sys.modules["safeguard"] = _real_sg_mod
            else: sys.modules.pop("safeguard", None)
            _ingest_mod._INGEST = {"proc": None, "path": "", "started": 0.0, "kind": None}
            os.environ.pop("VIEWER_INGEST_ROOTS", None)

        # =====================================================================================
        # do_POST's route-specific body-size cap (viewer_app.py): /api/ingest_upload gets a much
        # larger raw-body ceiling than every other POST route -- a real dragged PDF, base64-
        # encoded, is far bigger than the 8 MB MAX_POST_BYTES was ever sized for. Prove BOTH
        # halves with a real ~8.5 MB body over the actual socket (not mocked): the exception
        # route accepts it (doesn't 413 -- it may still fail validation for other reasons, e.g.
        # not a real PDF, but must get PAST the length check to do so) and every other route's
        # original cap is unweakened.
        # =====================================================================================
        import http.client as _httpclient

        declared = 8500000   # ~8.5 MB -- over the old 8 MB cap, comfortably under the new 200 MB one

        # /api/ingest (no upload exception): do_POST rejects purely on the Content-Length HEADER,
        # before ever calling self.rfile.read() -- so the body genuinely never needs to be sent to
        # prove the rejection, and skipping it sidesteps a client/server framing race (the server
        # closing the connection while a client is still mid-write of a real oversized body can
        # surface as a raw socket error -- e.g. WinError 10053 on Windows -- rather than a clean
        # HTTP response; that's an artifact of the test transport, not the feature under test).
        conn = _httpclient.HTTPConnection("127.0.0.1", PORT, timeout=15)
        try:
            conn.putrequest("POST", "/api/ingest")
            conn.putheader("Content-Type", "application/json")
            conn.putheader("Content-Length", str(declared))
            conn.endheaders()
            r = conn.getresponse()
            check("POST /api/ingest (a route WITHOUT the upload exception) still rejects an oversized "
                  "Content-Length -> 413 (cap not weakened globally)", r.status == 413)
        finally:
            conn.close()

        # /api/ingest_upload: genuinely under its 200 MB cap, so do_POST proceeds to read the full
        # body -- send it for real (the accept path has no early-close, so no race to sidestep).
        big_body = json.dumps({"filename": "big.pdf", "data": "A" * declared}).encode()
        req = urllib.request.Request(BASE + "/api/ingest_upload", data=big_body, headers={"Content-Type": "application/json"})
        try:
            with urllib.request.urlopen(req, timeout=20) as x:
                c = x.status
        except urllib.error.HTTPError as e:
            c = e.code
        check("POST /api/ingest_upload accepts a body well over the OLD 8 MB cap (not 413)", c != 413)

        # =====================================================================================
        # Genuine, UNMOCKED end-to-end upload: real HTTP POST with a real tiny PDF -> real
        # subprocess -> real crawl -> the uploaded document actually lands in the DB. Same split
        # as the progress-stamping e2e section below (validation is mocked above; this proves the
        # whole thing actually works for a real user, start to finish).
        #
        # Uses a SEPARATE, genuinely fresh + fully-migrated DB for this one check -- NOT the
        # shared fixture `db` the rest of this file uses -- and points the live server at it only
        # for the duration of this check (restored after). fixture.py hand-builds its tables to
        # match the CURRENT schema shape directly, without ever stamping schema_meta/schema_version
        # to say "every migration is already applied" -- fine for every other test here (they never
        # call migrate() against it), but the REAL subprocess this check launches always calls
        # viewer_ingest.py's real migrate() on startup, which would then try to re-apply migrations
        # whose columns the fixture already has, e.g. "ALTER TABLE pages ADD COLUMN ocr_priority" ->
        # "duplicate column name" (caught live: the subprocess exited 1, silently, since _launch()
        # redirects its stdout/stderr to DEVNULL -- reproduced directly by running the same CLI
        # command by hand to see the real traceback).
        # =====================================================================================
        try:
            import pymupdf as _fitz3
        except Exception:
            _fitz3 = None
        if _fitz3 is None:
            print("SKIP real end-to-end upload check (PyMuPDF not installed)")
        else:
            import viewer_ingest as _VI_upload
            up_dir = tempfile.mkdtemp(prefix="ingest_upload_e2e_")
            up_db = os.path.join(up_dir, "viewer.db")
            up_con = _VI_upload.connect(up_db)
            _VI_upload.migrate(up_con, os.path.join(ENGINE, "migrations"), db_path=up_db)
            up_con.close()

            doc3 = _fitz3.open(); pg3 = doc3.new_page(width=612, height=792)
            pg3.insert_text((72, 72), "REAL UPLOADED PDF FOR END-TO-END ROUTE TEST NSN 5305-01-999-8888")
            up_bytes = doc3.write(); doc3.close()
            import base64 as _b64e
            up_payload = {"filename": "REAL-UPLOAD-E2E.pdf", "data": _b64e.b64encode(up_bytes).decode()}

            orig_db_path = V.DB_PATH
            try:
                V.DB_PATH = up_db
                c, b = _req("/api/ingest_upload", up_payload)
                r = _json(b)
                check("real e2e upload -> ok + started (unmocked, genuinely launches viewer_ingest.py)",
                      c == 200 and r.get("ok") is True and r.get("started") is True)
                deadline = time.time() + 30
                landed = False
                while time.time() < deadline:
                    cs, bs = _req("/api/ingest_status")
                    st = _json(bs)
                    if not st.get("running"):
                        con_check = sqlite3.connect(up_db)
                        landed = con_check.execute(
                            "SELECT COUNT(*) FROM documents WHERE path LIKE ?", ("%REAL-UPLOAD-E2E.pdf",)
                        ).fetchone()[0] > 0
                        con_check.close()
                        break
                    time.sleep(0.5)
                check("real e2e upload: the uploaded document actually landed in the DB (real subprocess ran to completion)", landed)
            finally:
                V.DB_PATH = orig_db_path

        # =====================================================================================
        # GET /api/ingest_status -- richer response shape: 'progress' (live per-item detail, read
        # off the ingest_progress.json sidecar) and 'ocr_pending' (drives the UI's backlog card),
        # both additive -- neither breaks the pre-existing 'running'/'path'/'run' fields.
        # =====================================================================================
        c, b = _req("/api/ingest_status")
        r = _json(b)
        check("ingest_status -> 200", c == 200)
        check("ingest_status still has the original 'running' field", "running" in r)
        check("ingest_status has 'progress' (None when no sidecar written yet)", "progress" in r and r.get("progress") is None)
        check("ingest_status has 'ocr_pending' as an int", isinstance(r.get("ocr_pending"), int))

        prog_path = os.path.join(os.path.dirname(db), "ingest_progress.json")
        try:
            with open(prog_path, "w") as f:
                json.dump({"stage": "ocr", "current": {"doc": "TM-TEST.pdf", "page": 3}, "done": 7, "fail": 0, "total": 20}, f)
            c, b = _req("/api/ingest_status")
            r = _json(b)
            prog = r.get("progress") or {}
            check("ingest_status surfaces a real progress sidecar's stage", prog.get("stage") == "ocr")
            check("ingest_status surfaces the live 'current' doc/page detail",
                  (prog.get("current") or {}).get("doc") == "TM-TEST.pdf" and (prog.get("current") or {}).get("page") == 3)
            check("ingest_status surfaces done/total counts", prog.get("done") == 7 and prog.get("total") == 20)
        finally:
            try: os.unlink(prog_path)
            except OSError: pass

        # =====================================================================================
        # POST /api/airgap_manifest -- real signed manifest for a real folder + fence enforcement
        # =====================================================================================
        import airgap
        am_dir = tempfile.mkdtemp(prefix="airgap_manifest_src_")
        open(os.path.join(am_dir, "TM-A.pdf"), "wb").write(b"%PDF-1.4 alpha" + b"\x00" * 300)
        open(os.path.join(am_dir, "TM-B.pdf"), "wb").write(b"%PDF-1.4 bravo" + b"\x01" * 300)
        secret = "route-test-secret-2026"

        c, b = _req("/api/airgap_manifest", {"folder": am_dir, "secret": secret})
        r = _json(b)
        check("airgap_manifest -> 200 ok", c == 200 and r.get("ok") is True)
        manifest = r.get("manifest") or {}
        check("airgap_manifest counts both files", manifest.get("count") == 2)
        check("airgap_manifest signature validates with the real secret", airgap.signature_valid(manifest, secret))
        check("airgap_manifest signature rejects the wrong secret", not airgap.signature_valid(manifest, "wrong-secret"))

        fence_root3 = tempfile.mkdtemp(prefix="airgap_manifest_fenceroot_")
        os.environ["VIEWER_INGEST_ROOTS"] = fence_root3
        try:
            c, b = _req("/api/airgap_manifest", {"folder": am_dir, "secret": secret})
            check("airgap_manifest fence rejects a folder outside VIEWER_INGEST_ROOTS -> 400", c == 400)
        finally:
            os.environ.pop("VIEWER_INGEST_ROOTS", None)

        # =====================================================================================
        # POST /api/airgap_verify -- accept / tamper-reject / wrong-secret-reject + fence
        # =====================================================================================
        av_src = tempfile.mkdtemp(prefix="airgap_verify_src_")
        open(os.path.join(av_src, "TM-C.pdf"), "wb").write(b"%PDF-1.4 charlie" + b"\x02" * 300)
        av_secret = "route-verify-secret-2026"
        av_manifest = airgap.make_manifest(av_src, ["TM-C.pdf"], av_secret)

        av_dst = tempfile.mkdtemp(prefix="airgap_verify_dst_")
        shutil.copy(os.path.join(av_src, "TM-C.pdf"), os.path.join(av_dst, "TM-C.pdf"))

        c, b = _req("/api/airgap_verify", {"manifest": av_manifest, "folder": av_dst, "secret": av_secret})
        r = _json(b)
        result = r.get("result") or {}
        check("airgap_verify -> 200 ok", c == 200 and r.get("ok") is True)
        check("airgap_verify accepts a clean transfer", result.get("ok") is True and result.get("verdict") == "ACCEPT")

        with open(os.path.join(av_dst, "TM-C.pdf"), "r+b") as f:
            f.seek(20); f.write(b"\xff\xff")
        c, b = _req("/api/airgap_verify", {"manifest": av_manifest, "folder": av_dst, "secret": av_secret})
        r = _json(b); result = r.get("result") or {}
        check("airgap_verify rejects a tampered file", result.get("ok") is False and result.get("verdict") == "REJECT")
        check("airgap_verify reports the tampered filename", "TM-C.pdf" in (result.get("tampered") or []))

        shutil.copy(os.path.join(av_src, "TM-C.pdf"), os.path.join(av_dst, "TM-C.pdf"))   # restore clean copy
        c, b = _req("/api/airgap_verify", {"manifest": av_manifest, "folder": av_dst, "secret": "attacker-guess"})
        r = _json(b); result = r.get("result") or {}
        check("airgap_verify rejects the wrong secret", result.get("ok") is False and result.get("signature_valid") is False)

        fence_root4 = tempfile.mkdtemp(prefix="airgap_verify_fenceroot_")
        os.environ["VIEWER_INGEST_ROOTS"] = fence_root4
        try:
            c, b = _req("/api/airgap_verify", {"manifest": av_manifest, "folder": av_dst, "secret": av_secret})
            check("airgap_verify fence rejects a folder outside VIEWER_INGEST_ROOTS -> 400", c == 400)
        finally:
            os.environ.pop("VIEWER_INGEST_ROOTS", None)

        # =====================================================================================
        # POST /api/form_2404 + /api/form_2407 -- filled payload data actually renders into the PDF
        # =====================================================================================
        try:
            import forms as _forms
        except Exception:
            _forms = None
        if _forms is not None and _forms.available():
            _fitz = None
            for _modname in ("pymupdf", "fitz"):
                try:
                    _fitz = __import__(_modname); break
                except Exception:
                    continue
            if _fitz is not None:
                payload_2404 = {
                    "equipment": {"admin_no": "RT-2404-ADMIN-99", "nomenclature": "ROUTE TEST TRUCK"},
                    "faults": [{"item": "5", "deficiency": "ROUTE_TEST_DEFICIENCY_TEXT",
                                "corrective": "ROUTE_TEST_CORRECTIVE_TEXT", "status": "X"}],
                }
                c, b = _req("/api/form_2404", payload_2404)
                check("POST /api/form_2404 -> 200 PDF", c == 200 and b[:5] == b"%PDF-")
                doc = _fitz.open(stream=b, filetype="pdf")
                text_2404 = "\n".join(p.get_text() for p in doc)
                doc.close()
                check("POST /api/form_2404 fills admin_no into the PDF", "RT-2404-ADMIN-99" in text_2404)
                check("POST /api/form_2404 fills deficiency text into the PDF", "ROUTE_TEST_DEFICIENCY_TEXT" in text_2404)
                check("POST /api/form_2404 fills corrective text into the PDF", "ROUTE_TEST_CORRECTIVE_TEXT" in text_2404)

                payload_2407 = {
                    "organization": "RT-2407-ORG", "wo_no": "RT-WO-0099",
                    "equipment": {"admin_no": "RT-2407-ADMIN-77", "nomenclature": "ROUTE TEST TRAILER"},
                    "fault": "ROUTE_TEST_FAULT_NARRATIVE", "work_requested": "ROUTE_TEST_WORK_REQUESTED_TEXT",
                }
                c, b = _req("/api/form_2407", payload_2407)
                check("POST /api/form_2407 -> 200 PDF", c == 200 and b[:5] == b"%PDF-")
                doc = _fitz.open(stream=b, filetype="pdf")
                text_2407 = "\n".join(p.get_text() for p in doc)
                doc.close()
                check("POST /api/form_2407 fills admin_no into the PDF", "RT-2407-ADMIN-77" in text_2407)
                check("POST /api/form_2407 fills fault text into the PDF", "ROUTE_TEST_FAULT_NARRATIVE" in text_2407)
                check("POST /api/form_2407 fills work_requested text into the PDF", "ROUTE_TEST_WORK_REQUESTED_TEXT" in text_2407)
            else:
                print("SKIP form_2404/2407 filled-payload text checks (PyMuPDF not installed)")
        else:
            print("SKIP form_2404/2407 filled-payload text checks (reportlab not installed)")

        # =====================================================================================
        # End-to-end progress stamping: viewer_ingest.py's crawl() -> ocr() -> extract_parts()
        # each stamp ingest_progress.json as they actually run, against a REAL tiny PDF corpus --
        # not mocked. Deliberately engine-independent: ocr()'s _write_progress() call sits inside
        # handle(), which runs on BOTH the success and failure path (see viewer_ingest.py's OCR-
        # engine-failure fix earlier this session), so this holds whether or not tesseract/RapidOCR
        # is actually installed here -- unlike test_barcode_wiring.py's pipeline sections, nothing
        # here needs OCR to actually SUCCEED, only to actually RUN and reach handle() either way.
        # =====================================================================================
        try:
            import pymupdf as _fitz2
        except Exception:
            _fitz2 = None
        if _fitz2 is None:
            print("SKIP ingest_progress.json e2e stamping checks (PyMuPDF not installed)")
        else:
            import viewer_ingest as _VI
            e2e_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_")
            e2e_db = os.path.join(e2e_dir, "viewer.db")
            e2e_con = _VI.connect(e2e_db)
            _VI.migrate(e2e_con, os.path.join(ENGINE, "migrations"), db_path=e2e_db)
            corpus_dir = os.path.join(e2e_dir, "corpus"); os.makedirs(corpus_dir)
            prog_e2e_path = os.path.join(e2e_dir, "ingest_progress.json")

            # a text-layer page (indexed directly by crawl(), never queued for OCR) + an image-only
            # page (queued -- exercises the 'ocr' stage regardless of whether OCR actually succeeds).
            doc1 = _fitz2.open(); p1 = doc1.new_page(width=612, height=792)
            p1.insert_text((72, 72), "REAL TEXT LAYER PAGE FOR PROGRESS E2E TEST")
            doc1.save(os.path.join(corpus_dir, "TEXT-DOC.pdf")); doc1.close()

            doc2 = _fitz2.open(); p2 = doc2.new_page(width=612, height=792)
            # a filled black rectangle -- no text layer, dense enough to clear index_pdf()'s
            # blank-page skip so it genuinely gets queued for OCR (same shape test_barcode_wiring.py's
            # fixtures use, minus the barcode).
            p2.draw_rect(_fitz2.Rect(50, 50, 550, 740), color=(0, 0, 0), fill=(0, 0, 0))
            doc2.save(os.path.join(corpus_dir, "SCAN-DOC.pdf")); doc2.close()

            _VI.crawl(e2e_con, corpus_dir)
            e2e_con.commit()
            check("e2e crawl: text-layer page indexed directly (0 queued from it)",
                  e2e_con.execute("SELECT char_count FROM pages p JOIN documents d ON d.id=p.document_id "
                                  "WHERE d.path LIKE ?", ("%TEXT-DOC.pdf",)).fetchone()[0] > 0)
            queued = e2e_con.execute("SELECT COUNT(*) FROM pages WHERE ocr_status='pending'").fetchone()[0]
            check("e2e crawl: the image-only page was queued for OCR", queued == 1)
            with open(prog_e2e_path) as f:
                prog_after_crawl = json.load(f)
            check("e2e crawl stamped stage='crawl' in ingest_progress.json", prog_after_crawl.get("stage") == "crawl")
            check("e2e crawl stamped a real 'current' filename (not left over from init)",
                  isinstance(prog_after_crawl.get("current"), str) and prog_after_crawl["current"].endswith(".pdf"))
            # 'extracted' -- the "where is my data going" breakdown tally the in-app scan UI reads
            # to show more than a stage bar: which document(s) were found, and their metadata, as
            # soon as crawl() actually determines it (not just "N files scanned").
            extr_crawl = prog_after_crawl.get("extracted") or {}
            crawl_docs = extr_crawl.get("documents") or []
            check("e2e crawl tally: both documents recorded with a real db id", len(crawl_docs) == 2
                  and all(d.get("id") for d in crawl_docs))
            check("e2e crawl tally: pages_text reflects the ONE text-layer page (not the image-only one)",
                  extr_crawl.get("pages_text") == 1)

            remaining = _VI.ocr(e2e_con, 10, workers=1)
            with open(prog_e2e_path) as f:
                prog_after_ocr = json.load(f)
            check("e2e ocr stamped stage='ocr'", prog_after_ocr.get("stage") == "ocr")
            check("e2e ocr stamped total=1 (one page was queued)", prog_after_ocr.get("total") == 1)
            check("e2e ocr stamped done+fail covering the one queued page",
                  (prog_after_ocr.get("done") or 0) + (prog_after_ocr.get("fail") or 0) == 1)
            cur = prog_after_ocr.get("current") or {}
            check("e2e ocr stamped the real 'current' doc/page it just processed",
                  cur.get("doc") == "SCAN-DOC.pdf" and cur.get("page") == 1)
            check("e2e ocr: no pages left pending regardless of OCR success/failure (each got a terminal status)",
                  remaining == 0)
            extr_ocr = prog_after_ocr.get("extracted") or {}
            check("e2e ocr tally: the crawl-stage document list survives into the ocr stage (same running total, not reset)",
                  len(extr_ocr.get("documents") or []) == 2)
            check("e2e ocr tally: pages_ocr_done + pages_ocr_fail together cover the one queued page",
                  (extr_ocr.get("pages_ocr_done") or 0) + (extr_ocr.get("pages_ocr_fail") or 0) == 1)

            _VI.extract_parts(e2e_con)
            with open(prog_e2e_path) as f:
                prog_after_parts = json.load(f)
            check("e2e extract_parts stamped stage='parts'", prog_after_parts.get("stage") == "parts")
            e2e_con.close()

            # =================================================================================
            # A SECOND e2e run, with real RPSTL-formatted content this time, specifically to
            # exercise the tally fields the above fixture can't (neither of its docs has any
            # parts-list text): parts_page/parts_barcode counts and nsn_samples, which is what
            # the breakdown panel's "View part ->" link is actually built from.
            # =================================================================================
            e2e2_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_parts_")
            e2e2_db = os.path.join(e2e2_dir, "viewer.db")
            e2e2_con = _VI.connect(e2e2_db)
            _VI.migrate(e2e2_con, os.path.join(ENGINE, "migrations"), db_path=e2e2_db)
            corpus2_dir = os.path.join(e2e2_dir, "corpus"); os.makedirs(corpus2_dir)
            prog2_path = os.path.join(e2e2_dir, "ingest_progress.json")

            doc3 = _fitz2.open(); p3 = doc3.new_page(width=612, height=792)
            p3.insert_text((72, 72),
                "FIG 4 TEST ASSEMBLY\nITEM NO PART NUMBER FSCM DESCRIPTION USABLE ON CODE QTY\n"
                "ITEM 1 NSN 5330-01-654-9999 GASKET SET QTY 1", fontsize=10)
            doc3.save(os.path.join(corpus2_dir, "RPSTL-DOC.pdf")); doc3.close()

            _VI.crawl(e2e2_con, corpus2_dir); e2e2_con.commit()
            _VI.extract_parts(e2e2_con)
            with open(prog2_path) as f:
                prog2 = json.load(f)
            extr2 = prog2.get("extracted") or {}
            check("e2e parts tally: parts_page counts the real extracted NSN", extr2.get("parts_page") == 1)
            check("e2e parts tally: parts_barcode is 0 (no barcode on this page)", extr2.get("parts_barcode") == 0)
            check("e2e parts tally: nsn_samples contains the actual extracted NSN (what the UI links to)",
                  extr2.get("nsn_samples") == ["5330-01-654-9999"])
            e2e2_con.close()

            # =================================================================================
            # A THIRD e2e run: dimensional-data (measures.db) + schematic detection (both the
            # vector-native and raster/keyword paths), through the REAL 'run' subcommand's exact
            # stage sequence (crawl -> ocr -> extract_parts -> _run_schematic_stage), not just the
            # individual functions called in isolation like the sections above.
            # =================================================================================
            _VI._tally_reset()   # main() does this once per subprocess invocation; do the same here
                                 # since this whole test file shares one process across many e2e runs
            e2e3_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_dims_schem_")
            e2e3_db = os.path.join(e2e3_dir, "viewer.db")
            e2e3_con = _VI.connect(e2e3_db)
            _VI.migrate(e2e3_con, os.path.join(ENGINE, "migrations"), db_path=e2e3_db)
            corpus3_dir = os.path.join(e2e3_dir, "corpus"); os.makedirs(corpus3_dir)
            prog3_path = os.path.join(e2e3_dir, "ingest_progress.json")

            # a real-dimensioned text-layer page
            doc4 = _fitz2.open(); p4 = doc4.new_page(width=612, height=792)
            p4.insert_text((72, 72), "Bushing length 2.5 in, diameter .500 +/- .002 in. Thread 1/2-13 UNC-2A.")
            doc4.save(os.path.join(corpus3_dir, "DIMS-DOC.pdf")); doc4.close()

            # a vector-native schematic page: enough drawn lines to clear both has_vector (>=12
            # paths) and the >=8-edge min-edges floor _run_schematic_stage() applies.
            doc5 = _fitz2.open(); p5 = doc5.new_page(width=612, height=792)
            for i in range(8):
                y = 60 + i * 20
                p5.draw_line((60, y), (500, y)); p5.draw_line((60 + i * 10, 40), (60 + i * 10, 700))
            doc5.save(os.path.join(corpus3_dir, "VECTOR-SCHEM.pdf")); doc5.close()

            # a scanned/raster schematic page -- no vector content, caption-only detection.
            doc6 = _fitz2.open(); p6 = doc6.new_page(width=612, height=792)
            p6.draw_rect(_fitz2.Rect(40, 40, 570, 750), color=(0, 0, 0), fill=(0.9, 0.9, 0.9))
            p6.insert_text((80, 80), "FIG 14 WIRING DIAGRAM - MAIN HARNESS", fontsize=14)
            doc6.save(os.path.join(corpus3_dir, "RASTER-SCHEM.pdf")); doc6.close()

            _VI.crawl(e2e3_con, corpus3_dir)
            e2e3_con.commit()
            _VI.ocr(e2e3_con, 10, workers=1)
            _VI.extract_parts(e2e3_con)
            _VI._run_schematic_stage(e2e3_con)
            with open(prog3_path) as f:
                prog3 = json.load(f)
            extr3 = prog3.get("extracted") or {}

            meas_rows = sqlite3.connect(os.path.join(e2e3_dir, "measures.db")).execute(
                "SELECT type, value FROM meas").fetchall()
            check("e2e dimensions: real values landed in measures.db (not just an empty file)",
                  len(meas_rows) > 0 and any(v == "2.5" for _, v in meas_rows))
            check("e2e dimensions tally: 'dimensions' count matches what actually got written",
                  extr3.get("dimensions") == len(meas_rows))

            schem_rows = e2e3_con.execute(
                "SELECT document_id, detected_via, has_netlist, caption FROM schematics ORDER BY detected_via").fetchall()
            check("e2e schematics: exactly the two real schematic pages were detected (not the plain dims page)",
                  len(schem_rows) == 2)
            by_via = {r[1]: r for r in schem_rows}
            check("e2e schematics: the vector-native page got a real cached netlist", "vector" in by_via
                  and by_via["vector"][2] == 1)
            check("e2e schematics: the scanned page was caught by the keyword/caption path (no vector data to read)",
                  "keyword" in by_via and by_via["keyword"][3] and "WIRING DIAGRAM" in by_via["keyword"][3])
            check("e2e schematics tally: 'schematics' count matches what actually got written",
                  extr3.get("schematics") == len(schem_rows))
            schemcache_dir = os.path.join(e2e3_dir, "schemcache")
            check("e2e schematics: the netlist JSON sidecar exists at the SAME location BUILD-SCHEMGRAPH.bat "
                  "already uses (so /schemflow.js and Circuit Lab pick it up with no changes of their own)",
                  os.path.isdir(schemcache_dir) and len(os.listdir(schemcache_dir)) >= 1)

            # schematics_list() (features/browse_feature.py) must now surface a document purely
            # because it HAS a page-level schematics row -- even though neither VECTOR-SCHEM.pdf nor
            # RASTER-SCHEM.pdf's filename/title matches the old LIKE pattern at all.
            import features.browse_feature as _bf
            _orig_core = _bf.core
            class _FakeCoreForSchemList:
                DB_PATH = e2e3_db
                @staticmethod
                def db():
                    c = sqlite3.connect(_FakeCoreForSchemList.DB_PATH); c.row_factory = sqlite3.Row; return c
            _bf.core = _FakeCoreForSchemList
            try:
                slist = _bf.schematics_list()
                names = {it["filename"] for it in slist["items"]}
                check("schematics_list() now includes a document found ONLY via page-level detection "
                      "(filename never matched the old LIKE pattern)", "VECTOR-SCHEM.pdf" in names)
            finally:
                _bf.core = _orig_core

            e2e3_con.close()

            # =================================================================================
            # A FOURTH e2e run: the Discovery Engine's non-PDF format support (image/txt/html/
            # unsupported-office) plus table extraction, through the same real 'run' stage
            # sequence. Before this, classify_ext() could already TELL a .jpg from a .docx but
            # crawl() did nothing with either -- confirmed by grep, no code path existed at all.
            # =================================================================================
            try:
                from PIL import Image, ImageDraw, ImageFont
                _have_pil_font = True
            except Exception:
                _have_pil_font = False
            if not _have_pil_font:
                print("SKIP non-PDF format + tables e2e checks (PIL not installed)")
            else:
                e2e4_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_formats_")
                e2e4_db = os.path.join(e2e4_dir, "viewer.db")
                e2e4_con = _VI.connect(e2e4_db)
                _VI.migrate(e2e4_con, os.path.join(ENGINE, "migrations"), db_path=e2e4_db)
                corpus4_dir = os.path.join(e2e4_dir, "corpus"); os.makedirs(corpus4_dir)

                # a real, OCR-able standalone image (needs a real font -- the default PIL bitmap
                # font is too thin to clear the blank-page density skip, confirmed live earlier).
                try:
                    font = ImageFont.truetype("arial.ttf", 30)
                except Exception:
                    font = ImageFont.load_default()
                img = Image.new("RGB", (800, 200), "white")
                ImageDraw.Draw(img).text((30, 30), "BOLT LENGTH 2.0 IN NSN 5305-01-888-7777", fill="black", font=font)
                img.save(os.path.join(corpus4_dir, "photo.png"))

                with open(os.path.join(corpus4_dir, "notes.txt"), "w") as f:
                    f.write("Gasket length 3.0 in, diameter .750 in.")
                with open(os.path.join(corpus4_dir, "export.html"), "w") as f:
                    f.write("<html><body><h1>Torque Table</h1><p>Bolt torque 45 ft-lb <b>required</b>.</p></body></html>")
                with open(os.path.join(corpus4_dir, "manual.doc"), "wb") as f:
                    f.write(b"\xd0\xcf\x11\xe0 fake legacy-binary-Word bytes -- unsupported, must not crash the crawl")

                # a real ruled-line table (tables.py needs actual grid geometry, not just text --
                # confirmed live: a plain textbox never triggers PyMuPDF's find_tables()).
                doc7 = _fitz2.open(); p7 = doc7.new_page(width=612, height=792)
                x0, y0, cw, rh = 60, 60, 150, 30
                for r in range(4): p7.draw_line((x0, y0 + r * rh), (x0 + 3 * cw, y0 + r * rh))
                for c in range(4): p7.draw_line((x0 + c * cw, y0), (x0 + c * cw, y0 + 3 * rh))
                for r, row in enumerate([["NSN", "QTY", "TORQUE"], ["5305-01-111-2222", "4", "45 ft-lb"], ["5305-01-333-4444", "2", "30 ft-lb"]]):
                    for c, val in enumerate(row):
                        p7.insert_text((x0 + c * cw + 5, y0 + r * rh + 20), val, fontsize=10)
                doc7.save(os.path.join(corpus4_dir, "TABLE-DOC.pdf")); doc7.close()

                _VI.crawl(e2e4_con, corpus4_dir)
                e2e4_con.commit()

                doc_rows = {os.path.basename(p): (t, pc, st) for p, t, pc, st in
                           e2e4_con.execute("SELECT path, type, page_count, status FROM documents")}
                check("e2e formats: image classified + given exactly 1 pending page (queued for OCR)",
                      doc_rows.get("photo.png") == ("image", 1, "partial"))
                check("e2e formats: .txt read directly, indexed immediately (no OCR needed)",
                      doc_rows.get("notes.txt") == ("text", 1, "indexed"))
                check("e2e formats: .html tag-stripped + indexed immediately",
                      doc_rows.get("export.html") == ("html", 1, "indexed"))
                check("e2e formats: unsupported .doc (legacy binary) discovered but not crashed, 0 pages",
                      doc_rows.get("manual.doc") == ("office", 0, "indexed"))

                txt_body = e2e4_con.execute(
                    "SELECT body_text FROM pages p JOIN documents d ON d.id=p.document_id WHERE d.path LIKE ?",
                    ("%notes.txt",)).fetchone()[0]
                check("e2e formats: .txt body_text is the real file content", "Gasket length 3.0 in" in txt_body)
                html_body = e2e4_con.execute(
                    "SELECT body_text FROM pages p JOIN documents d ON d.id=p.document_id WHERE d.path LIKE ?",
                    ("%export.html",)).fetchone()[0]
                check("e2e formats: .html tags stripped, real text survives, markup does not",
                      "Torque Table" in html_body and "<h1>" not in html_body and "<b>" not in html_body)

                remaining4 = _VI.ocr(e2e4_con, 10, workers=1)
                check("e2e formats: the standalone image actually got OCR'd (no pages left pending)", remaining4 == 0)
                img_row = e2e4_con.execute(
                    "SELECT body_text, ocr_status FROM pages p JOIN documents d ON d.id=p.document_id WHERE d.path LIKE ?",
                    ("%photo.png",)).fetchone()
                check("e2e formats: the image's OCR'd text is real (the SAME OCR pipeline a scanned PDF page uses)",
                      img_row[1] == "done" and "BOLT LENGTH" in img_row[0])

                _VI.extract_parts(e2e4_con)
                _VI._run_schematic_stage(e2e4_con)
                _VI._run_tables_stage(e2e4_con)

                meas4 = sqlite3.connect(os.path.join(e2e4_dir, "measures.db")).execute(
                    "SELECT value FROM meas WHERE type='length'").fetchall()
                check("e2e formats: dimensional extraction ALSO ran on the OCR'd image and the .txt/.html pages "
                      "(not just PDFs)", len(meas4) >= 3)   # image "2.0", txt "3.0", html has no length but torque

                tbl4 = sqlite3.connect(os.path.join(e2e4_dir, "tables.db")).execute(
                    "SELECT n_rows, n_cols, spec, units FROM tbl").fetchall()
                check("e2e tables: the real ruled table was detected with correct shape + spec/units flags",
                      tbl4 == [(3, 3, 1, "torque")])

                e2e4_con.close()
                # (ingest_upload()'s own accept/reject behavior for these same formats is already
                # covered end-to-end by the mocked-Popen section above -- .docx rejection, .txt
                # acceptance, and bad-image-bytes rejection -- so it isn't re-tested here.)

            # =================================================================================
            # A FIFTH e2e run: the full-codebase audit's three genuine wire-ins -- RPSTL parts-list
            # row extraction (same "built, real live consumer routes, only ever ran via a manual
            # .bat" gap already fixed for measures/schematics/tables), header/footer boilerplate
            # stripping (pagetrim.py) on text-layer pages before they're stored/measured, and the
            # enrich_flis() -> build_keywords.py hookup (keywords.json refresh right after the
            # colloquial names it depends on are populated, instead of a manual second step).
            # =================================================================================
            _VI._tally_reset()
            e2e5_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_audit_")
            e2e5_db = os.path.join(e2e5_dir, "viewer.db")
            e2e5_con = _VI.connect(e2e5_db)
            _VI.migrate(e2e5_con, os.path.join(ENGINE, "migrations"), db_path=e2e5_db)
            corpus5_dir = os.path.join(e2e5_dir, "corpus"); os.makedirs(corpus5_dir)

            # --- RPSTL: a real RPSTL-named doc with a real parts-list line, + a same-shaped line in
            # a doc whose name/tm_number never matches the RPSTL_LIKE gate (imported straight from
            # build_rpstl.py, not duplicated) -- must NOT get scanned. ---
            doc8 = _fitz2.open(); p8 = doc8.new_page(width=612, height=792)
            p8.insert_text((50, 50), "12  PAOZZ  5305-01-234-5678  81349  MS35206-243   SCREW,MACHINE   4")
            doc8.save(os.path.join(corpus5_dir, "TM-9-2320-000-24P.pdf")); doc8.close()
            doc9 = _fitz2.open(); p9 = doc9.new_page(width=612, height=792)
            p9.insert_text((50, 50), "12  PAOZZ  5305-01-999-0000  81349  MS35206-244   SCREW,MACHINE   2")
            doc9.save(os.path.join(corpus5_dir, "TM-9-2320-000-10.pdf")); doc9.close()   # operator manual, not RPSTL-like

            # --- pagetrim: a real >=5-page doc with a genuinely recurring header/footer AND real,
            # page-to-page-VARYING body content (pagetrim only flags text that recurs verbatim
            # across pages -- identical body content on every page would also get flagged, so this
            # mirrors pagetrim.py's own self-test: rotate distinct words per page). ---
            pt_header = "TM 9-8888-777-14"
            pt_words = ["alternator", "bracket", "coolant", "differential", "engine", "flywheel", "gasket", "harness", "injector", "manifold"]
            doc10 = _fitz2.open()
            for i in range(8):
                p = doc10.new_page(width=612, height=792)
                y = 50
                p.insert_text((50, y), pt_header); y += 20
                p.insert_text((50, y), "SECTION II MAINTENANCE"); y += 30
                for j in range(10):
                    w1 = pt_words[(i + j) % len(pt_words)]; w2 = pt_words[(i + j + 3) % len(pt_words)]
                    p.insert_text((50, y), "Install the %s and torque the %s fitting to %d ft-lb." % (w1, w2, 40 + j))
                    y += 18
                p.insert_text((50, 750), "Change 2   Page %d" % (12 + i))
            doc10.save(os.path.join(corpus5_dir, "TM-9-8888-777-14.pdf")); doc10.close()

            _VI.crawl(e2e5_con, corpus5_dir)
            e2e5_con.commit()
            _VI.extract_parts(e2e5_con)
            _VI._run_rpstl_stage(e2e5_con)

            rpstl_rows = sqlite3.connect(os.path.join(e2e5_dir, "rpstl.db")).execute(
                "SELECT pn_norm, nsn, smr, nomenclature FROM parts_rows").fetchall()
            check("e2e RPSTL: exactly one row extracted (the RPSTL-named doc, not the operator manual)",
                  len(rpstl_rows) == 1)
            if rpstl_rows:
                check("e2e RPSTL: the row's fields are real (NSN + SMR correctly parsed)",
                      rpstl_rows[0][1] == "5305-01-234-5678" and rpstl_rows[0][2] == "PAOZZ")
            check("e2e RPSTL tally: 'rpstl' count matches what actually got written",
                  _VI._EXTRACT_TALLY.get("rpstl") == len(rpstl_rows))
            _VI._run_rpstl_stage(e2e5_con)   # idempotent re-run must not duplicate rows
            rpstl_n2 = sqlite3.connect(os.path.join(e2e5_dir, "rpstl.db")).execute(
                "SELECT COUNT(*) FROM parts_rows").fetchone()[0]
            check("e2e RPSTL: idempotent re-run does not duplicate rows", rpstl_n2 == len(rpstl_rows))

            pt_doc_id = e2e5_con.execute("SELECT id FROM documents WHERE path LIKE ?", ("%TM-9-8888-777-14.pdf",)).fetchone()[0]
            pt_body1 = e2e5_con.execute("SELECT body_text FROM pages WHERE document_id=? AND page_number=1",
                                        (pt_doc_id,)).fetchone()[0]
            check("e2e pagetrim: the recurring header is stripped from the stored page body", pt_header not in pt_body1)
            check("e2e pagetrim: the recurring footer is stripped from the stored page body", "Change 2" not in pt_body1)
            check("e2e pagetrim: real, non-recurring body content survives", "Install the" in pt_body1 and "ft-lb" in pt_body1)
            pt_drow = e2e5_con.execute("SELECT tm_number FROM documents WHERE id=?", (pt_doc_id,)).fetchone()
            check("e2e pagetrim: tm_number metadata is STILL correctly detected from the RAW header text "
                  "(meta_text is computed before stripping, never from the cleaned copy)",
                  pt_drow[0] == pt_header)

            # toggle off -> the header must survive untouched (proves the stripping is real, not a no-op)
            e2e5b_db = os.path.join(e2e5_dir, "viewer_pt_off.db")
            e2e5b_con = _VI.connect(e2e5b_db)
            _VI.migrate(e2e5b_con, os.path.join(ENGINE, "migrations"), db_path=e2e5b_db)
            _VI.PAGETRIM_SCAN = False
            try:
                _VI._tally_reset()
                _VI.crawl(e2e5b_con, corpus5_dir)
                e2e5b_con.commit()
                pt_doc_id_off = e2e5b_con.execute("SELECT id FROM documents WHERE path LIKE ?",
                                                  ("%TM-9-8888-777-14.pdf",)).fetchone()[0]
                pt_body1_off = e2e5b_con.execute("SELECT body_text FROM pages WHERE document_id=? AND page_number=1",
                                                 (pt_doc_id_off,)).fetchone()[0]
                check("e2e pagetrim: PAGETRIM_SCAN=False -> the header is NOT stripped", pt_header in pt_body1_off)
                # flags.py's registry (added alongside this same audit) must reflect this live --
                # _write_progress() bakes flags_off into EVERY progress write automatically.
                prog_off = json.load(open(os.path.join(e2e5_dir, "ingest_progress.json")))
                check("e2e flags: ingest_progress.json's flags_off includes 'pagetrim' while it's toggled off",
                      "pagetrim" in (prog_off.get("flags_off") or []))
            finally:
                _VI.PAGETRIM_SCAN = True
                e2e5b_con.close()
            prog_back_on = json.load(open(os.path.join(e2e5_dir, "ingest_progress.json")))
            # (the file itself is a stale snapshot from the toggled-off run above -- re-derive fresh
            # from the registry directly, which is the actual live-vs-snapshot guarantee under test)
            import flags as _flags
            check("e2e flags: flags.disabled_stage_names() no longer includes 'pagetrim' once restored",
                  "pagetrim" not in _flags.disabled_stage_names())

            # `python viewer_ingest.py flags` -- the CLI introspection entry point, exercised as a
            # real subprocess (same as every other CLI-surface check in this file).
            flags_proc = subprocess.run(
                [sys.executable, os.path.join(ENGINE, "viewer_ingest.py"), "flags"],
                capture_output=True, text=True, timeout=30)
            check("`viewer_ingest.py flags` exits 0", flags_proc.returncode == 0)
            check("`viewer_ingest.py flags` lists all 9 real toggles by env var name", all(
                e in flags_proc.stdout for e in (
                    "VIEWER_OCR_PREPROCESS", "VIEWER_BARCODE_SCAN", "VIEWER_MEASURES_SCAN",
                    "VIEWER_SCHEMATIC_SCAN", "VIEWER_TABLES_SCAN", "VIEWER_RPSTL_SCAN",
                    "VIEWER_PAGETRIM_SCAN", "VIEWER_KEYWORDS_SCAN", "VIEWER_OFFICE_SCAN")))
            check("`viewer_ingest.py flags` reports all 9 active by default (no env override)",
                  "9 of 9 toggles active" in flags_proc.stdout)

            # --- keywords: enrich_flis() populating a real 'Also called:' colloquial name must
            # trigger build_keywords.run() -- monkeypatched to a recording stub so this NEVER writes
            # the real engine/keywords.json. ---
            rpstl_doc_id = e2e5_con.execute("SELECT id FROM documents WHERE path LIKE ?", ("%TM-9-2320-000-24P.pdf",)).fetchone()[0]
            e2e5_con.execute("UPDATE documents SET nsn=? WHERE id=?", ("5305-01-234-5678", rpstl_doc_id))
            e2e5_con.commit()

            flis_dir = os.path.join(e2e5_dir, "flis")
            os.makedirs(flis_dir)
            def _write_csv(name, rows):
                with open(os.path.join(flis_dir, name), "w", encoding="utf-8", newline="") as f:
                    f.write("HEADER_ROW_SKIPPED\n")
                    for r in rows: f.write(",".join(r) + "\n")
            _write_csv("V_FLIS_IDENTIFICATION.CSV", [("012345678", "12345")])          # NIIN,INC
            _write_csv("P_H6_PICK.CSV", [("12345", "SCREW MACHINE")])                  # INC,item name
            _write_csv("V_COLLOQUIAL_NAME.CSV", [("12345", "", "cap screw")])          # INC,RELATED_INC,colloquial

            import build_keywords as _bk
            _orig_bk_run = _bk.run
            _bk_calls = []
            _bk.run = lambda db=None, out=None: (_bk_calls.append(db), (1, 1, 0, True))[-1]
            try:
                n_enriched = _VI.enrich_flis(e2e5_con, flis_dir)
                check("e2e keywords wiring: enrich_flis() actually enriched >=1 NSN (fixture is valid)", n_enriched >= 1)
                check("e2e keywords wiring: enrich_flis() called build_keywords.run() exactly once after enriching",
                      len(_bk_calls) == 1)
                check("e2e keywords wiring: build_keywords.run() was called with the SAME db enrich_flis() just wrote to",
                      bool(_bk_calls) and _bk_calls[0] and
                      os.path.normcase(os.path.abspath(_bk_calls[0])) == os.path.normcase(os.path.abspath(e2e5_db)))
            finally:
                _bk.run = _orig_bk_run

            _bk_calls2 = []
            _bk.run = lambda db=None, out=None: (_bk_calls2.append(db), (0, 0, 0, True))[-1]
            try:
                empty_flis_dir = tempfile.mkdtemp(prefix="empty_flis_")
                n0 = _VI.enrich_flis(e2e5_con, empty_flis_dir)
                check("e2e keywords wiring: no NSNs enriched -> build_keywords.run() is NOT called",
                      n0 == 0 and len(_bk_calls2) == 0)
            finally:
                _bk.run = _orig_bk_run

            _bk_calls3 = []
            _bk.run = lambda db=None, out=None: (_bk_calls3.append(db), (0, 0, 0, True))[-1]
            _VI.KEYWORDS_SCAN = False
            try:
                n3 = _VI.enrich_flis(e2e5_con, flis_dir)
                check("e2e keywords wiring: KEYWORDS_SCAN=False -> build_keywords.run() is NOT called "
                      "even though NSNs WERE enriched", n3 >= 1 and len(_bk_calls3) == 0)
            finally:
                _VI.KEYWORDS_SCAN = True
                _bk.run = _orig_bk_run

            e2e5_con.close()

            # =================================================================================
            # A SIXTH e2e run: Office document text extraction (office.py -- .docx/.xlsx/.pptx/
            # .rtf), a deferred item picked up alongside the flags audit, tier-gated to the modern
            # OS signal per explicit user direction (avoid Vista/legacy incompatibilities).
            # =================================================================================
            try:
                import docx as _have_docx_mod
                import openpyxl as _have_xlsx_mod
                import pptx as _have_pptx_mod
                _have_office_libs = True
            except Exception:
                _have_office_libs = False
            if not _have_office_libs:
                print("SKIP Office-format e2e checks (python-docx/openpyxl/python-pptx not installed)")
            else:
                e2e6_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_office_")
                e2e6_db = os.path.join(e2e6_dir, "viewer.db")
                e2e6_con = _VI.connect(e2e6_db)
                _VI.migrate(e2e6_con, os.path.join(ENGINE, "migrations"), db_path=e2e6_db)
                corpus6_dir = os.path.join(e2e6_dir, "corpus"); os.makedirs(corpus6_dir)

                docx_doc = _have_docx_mod.Document()
                docx_doc.add_paragraph("Torque the alternator bracket bolt to 45 ft-lb.")
                docx_doc.add_paragraph("NSN 5305-01-123-4567 SCREW,MACHINE")
                docx_doc.save(os.path.join(corpus6_dir, "notes.docx"))

                wb = _have_xlsx_mod.Workbook()
                ws1 = wb.active; ws1.title = "Specs"
                ws1["A1"] = "Overall length"; ws1["B1"] = "180 in"
                wb.create_sheet("Torque")["A1"] = "45 ft-lb"
                wb.save(os.path.join(corpus6_dir, "specs.xlsx"))

                pr = _have_pptx_mod.Presentation()
                s1 = pr.slides.add_slide(pr.slide_layouts[1])
                s1.shapes.title.text = "Bolt Torque"
                s1.placeholders[1].text = "NSN 5305-01-123-4567 SCREW,MACHINE"
                pr.save(os.path.join(corpus6_dir, "deck.pptx"))

                with open(os.path.join(corpus6_dir, "memo.rtf"), "w", encoding="latin-1") as f:
                    f.write(r"{\rtf1\ansi{\fonttbl{\f0 Arial;}}\f0 Fording depth 30 in required.\par }")

                _VI._tally_reset()
                _VI.crawl(e2e6_con, corpus6_dir)
                e2e6_con.commit()

                doc6_rows = {os.path.basename(p): (t, pc, st) for p, t, pc, st in
                            e2e6_con.execute("SELECT path, type, page_count, status FROM documents")}
                check("e2e office: .docx classified + 1 page extracted", doc6_rows.get("notes.docx") == ("docx", 1, "indexed"))
                check("e2e office: .xlsx classified + 2 pages (one per sheet)", doc6_rows.get("specs.xlsx") == ("xlsx", 2, "indexed"))
                check("e2e office: .pptx classified + 1 page (one per slide)", doc6_rows.get("deck.pptx") == ("pptx", 1, "indexed"))
                check("e2e office: .rtf classified + 1 page (dependency-free, no tier gate)", doc6_rows.get("memo.rtf") == ("rtf", 1, "indexed"))

                docx_id = e2e6_con.execute("SELECT id FROM documents WHERE path LIKE ?", ("%notes.docx",)).fetchone()[0]
                docx_body = e2e6_con.execute("SELECT body_text FROM pages WHERE document_id=? AND page_number=1", (docx_id,)).fetchone()[0]
                check("e2e office: .docx body_text is the real paragraph content", "Torque the alternator bracket bolt" in docx_body)

                xlsx_id = e2e6_con.execute("SELECT id FROM documents WHERE path LIKE ?", ("%specs.xlsx",)).fetchone()[0]
                xlsx_pages = e2e6_con.execute("SELECT page_number, body_text FROM pages WHERE document_id=? ORDER BY page_number", (xlsx_id,)).fetchall()
                check("e2e office: .xlsx sheet 1 body_text has the real cell content", "180 in" in xlsx_pages[0][1])
                check("e2e office: .xlsx sheet 2 is a genuinely separate page", "45 ft-lb" in xlsx_pages[1][1] and "45 ft-lb" not in xlsx_pages[0][1])

                meas6 = sqlite3.connect(os.path.join(e2e6_dir, "measures.db")).execute(
                    "SELECT type, value FROM meas").fetchall()
                check("e2e office: dimensional extraction ran on Office-extracted text too "
                      "(docx torque, xlsx length+torque, rtf length)",
                      ("torque", "45") in meas6 and ("length", "180") in meas6 and ("length", "30") in meas6)

                # toggle off -> discovered, 0 pages, same degrade shape .doc/.xls/.ppt already have
                e2e6b_db = os.path.join(e2e6_dir, "viewer_off.db")
                e2e6b_con = _VI.connect(e2e6b_db)
                _VI.migrate(e2e6b_con, os.path.join(ENGINE, "migrations"), db_path=e2e6b_db)
                _VI.OFFICE_SCAN = False
                try:
                    _VI._tally_reset()
                    _VI.crawl(e2e6b_con, corpus6_dir)
                    e2e6b_con.commit()
                    off_row = e2e6b_con.execute("SELECT type, page_count, status FROM documents WHERE path LIKE ?",
                                                ("%notes.docx",)).fetchone()
                    check("e2e office: OFFICE_SCAN=False -> discovered, 0 pages (same shape as .doc/.xls/.ppt)",
                          off_row == ("office", 0, "indexed"))
                finally:
                    _VI.OFFICE_SCAN = True
                    e2e6b_con.close()

                # unsupported legacy binary formats (.doc/.xls/.ppt) still just discover, 0 pages, no crash
                with open(os.path.join(corpus6_dir, "legacy.doc"), "wb") as f:
                    f.write(b"\xd0\xcf\x11\xe0 fake legacy binary doc -- must not crash the crawl")
                _VI._tally_reset()
                _VI.crawl(e2e6_con, corpus6_dir)
                e2e6_con.commit()
                doc_row = e2e6_con.execute("SELECT type, page_count, status FROM documents WHERE path LIKE ?",
                                           ("%legacy.doc",)).fetchone()
                check("e2e office: .doc (legacy binary) still unsupported, discovered, 0 pages, no crash",
                      doc_row == ("office", 0, "indexed"))

                e2e6_con.close()

                # ingest_upload() must accept the same 4 formats, with real content validation
                # (ZIP/OOXML magic for docx/xlsx/pptx, {\rtf header for rtf). Self-contained imports
                # here rather than relying on _if/_b64u from the e2e4 section above, which are only
                # bound if PIL happened to be installed when that section ran. Popen/safeguard
                # mocked (reusing the fakes from the POST /api/ingest section above) -- two REAL,
                # unmocked launches back-to-back (the accepted .docx, then the accepted .rtf) would
                # otherwise collide with ingest_feature.py's own single-job-at-a-time lock
                # (_INGEST["proc"].poll() is None), rejecting the second as "already in progress"
                # for a reason that has nothing to do with the content validation under test here.
                import features.ingest_feature as _if6
                import base64 as _b64u
                _orig_core_if6 = _if6.core
                class _FakeCoreForOffice:
                    DB_PATH = e2e6_db
                _if6.core = _FakeCoreForOffice
                _subprocess_mod.Popen = _fake_popen
                sys.modules["safeguard"] = _FakeSafeguard()
                try:
                    real_docx_bytes = open(os.path.join(corpus6_dir, "notes.docx"), "rb").read()
                    r_docx = _if6.ingest_upload("upload.docx", _b64u.b64encode(real_docx_bytes).decode())
                    check("ingest_upload() accepts a real .docx (valid ZIP/OOXML header)", r_docx.get("ok") is True)

                    r_bad_docx = _if6.ingest_upload("fake.docx", _b64u.b64encode(b"not a real docx").decode())
                    check("ingest_upload() rejects .docx bytes that aren't a real ZIP/OOXML container",
                          r_bad_docx.get("ok") is False)

                    real_rtf_bytes = open(os.path.join(corpus6_dir, "memo.rtf"), "rb").read()
                    r_rtf = _if6.ingest_upload("upload.rtf", _b64u.b64encode(real_rtf_bytes).decode())
                    check("ingest_upload() accepts a real .rtf (valid {\\rtf header)", r_rtf.get("ok") is True)

                    r_bad_rtf = _if6.ingest_upload("fake.rtf", _b64u.b64encode(b"not rtf at all").decode())
                    check("ingest_upload() rejects .rtf bytes missing the {\\rtf header",
                          r_bad_rtf.get("ok") is False)

                    r_legacy_doc = _if6.ingest_upload("old.doc", _b64u.b64encode(b"\xd0\xcf\x11\xe0 legacy binary").decode())
                    check("ingest_upload() still rejects .doc (legacy binary, unsupported extension)",
                          r_legacy_doc.get("ok") is False)
                finally:
                    _if6.core = _orig_core_if6
                    _subprocess_mod.Popen = _real_popen
                    if _real_sg_mod is not None: sys.modules["safeguard"] = _real_sg_mod
                    else: sys.modules.pop("safeguard", None)

            # =================================================================================
            # A SEVENTH e2e run: edition/duplicate clustering (dedup.py + build_dedup.py + the
            # /api/editions route), another deferred item picked up alongside the flags audit.
            # build_dedup.py exercised as a REAL subprocess (same as viewer_ingest.py's `flags`
            # subcommand above) against a real corpus of near-duplicate + distinct documents.
            # =================================================================================
            e2e7_dir = tempfile.mkdtemp(prefix="ingest_progress_e2e_dedup_")
            e2e7_db = os.path.join(e2e7_dir, "viewer.db")
            e2e7_con = _VI.connect(e2e7_db)
            _VI.migrate(e2e7_con, os.path.join(ENGINE, "migrations"), db_path=e2e7_db)
            corpus7_dir = os.path.join(e2e7_dir, "corpus"); os.makedirs(corpus7_dir)

            base_text = ("The alternator is mounted on the front of the engine and is driven by the serpentine "
                         "belt. Remove the two mounting bolts and disconnect the wiring harness before extraction. "
                         "Torque to 30 foot pounds and verify clearance before reassembly of the housing.")
            edition_text = base_text.replace("30 foot pounds", "35 foot pounds") + " Change 3 revision notice page."
            other_text = ("The transmission fluid should be checked with the vehicle on level ground and the "
                          "engine at operating temperature. Use only the specified lubricant grade and do not "
                          "overfill the reservoir under any condition or void the warranty coverage terms.")

            doc11 = _fitz2.open(); p11 = doc11.new_page(width=612, height=792)
            p11.insert_text((72, 72), base_text)
            doc11.save(os.path.join(corpus7_dir, "TM-9-2320-280-24-CH2.pdf")); doc11.close()

            doc12 = _fitz2.open(); p12 = doc12.new_page(width=612, height=792)
            p12.insert_text((72, 72), edition_text)
            doc12.save(os.path.join(corpus7_dir, "TM-9-2320-280-24-CH3.pdf")); doc12.close()

            doc13 = _fitz2.open(); p13 = doc13.new_page(width=612, height=792)
            p13.insert_text((72, 72), other_text)
            doc13.save(os.path.join(corpus7_dir, "TM-9-2320-280-10.pdf")); doc13.close()

            _VI._tally_reset()
            _VI.crawl(e2e7_con, corpus7_dir)
            e2e7_con.commit()
            e2e7_con.close()

            e2e7_dedup_db = os.path.join(e2e7_dir, "dedup.db")
            env = dict(os.environ, VIEWER_DB=e2e7_db, DEDUP_DB=e2e7_dedup_db)
            dedup_proc = subprocess.run(
                [sys.executable, os.path.join(ENGINE, "build_dedup.py"), "--threshold", "0.6"],
                capture_output=True, text=True, timeout=60, env=env)
            check("build_dedup.py exits 0", dedup_proc.returncode == 0)
            check("build_dedup.py reports finding exactly one cluster", "1 cluster" in dedup_proc.stdout)
            check("build_dedup.py actually wrote dedup.db", os.path.exists(e2e7_dedup_db))

            e2e7_con2 = sqlite3.connect(e2e7_db)
            ch2_id = e2e7_con2.execute("SELECT id FROM documents WHERE path LIKE ?", ("%TM-9-2320-280-24-CH2.pdf",)).fetchone()[0]
            ch3_id = e2e7_con2.execute("SELECT id FROM documents WHERE path LIKE ?", ("%TM-9-2320-280-24-CH3.pdf",)).fetchone()[0]
            other_id = e2e7_con2.execute("SELECT id FROM documents WHERE path LIKE ?", ("%TM-9-2320-280-10.pdf",)).fetchone()[0]
            e2e7_con2.close()

            import dedup as _dedup_mod
            ch2_editions = _dedup_mod.editions_for(e2e7_dedup_db, ch2_id)
            check("e2e dedup: the near-duplicate CH2/CH3 pair was correctly clustered together",
                  len(ch2_editions) == 1 and ch2_editions[0]["document_id"] == ch3_id)
            check("e2e dedup: the genuinely different document has no editions", _dedup_mod.editions_for(e2e7_dedup_db, other_id) == [])

            # the real /api/editions route, exercised directly (same technique test_tables_plus_stitch.py
            # uses -- no need for a full HTTP server to prove the routing + real dedup.py calls are wired).
            import features.routes.doc_extractors as _de_mod
            class _FakeHandler7:
                def __init__(self): self.sent=None
                def _send(self, c, b): self.sent=(c,b)
            class _FakeCoreForEditions:
                DB_PATH = e2e7_db
            _orig_de_core = _de_mod.core
            _de_mod.core = _FakeCoreForEditions
            try:
                h7 = _FakeHandler7()
                _de_mod.r_editions(h7, {"doc": [str(ch2_id)]})
                code7, body7 = h7.sent
                check("/api/editions route: 200 + available=True (dedup.db exists)",
                      code7 == 200 and body7.get("available") is True)
                check("/api/editions route: returns the real sibling edition via the real route function",
                      len(body7.get("editions") or []) == 1 and body7["editions"][0]["document_id"] == ch3_id)
            finally:
                _de_mod.core = _orig_de_core
    finally:
        if _orig_roots is None: os.environ.pop("VIEWER_INGEST_ROOTS", None)
        else: os.environ["VIEWER_INGEST_ROOTS"] = _orig_roots
        srv.shutdown()

    fails = [n for n, ok in tests if not ok]
    for n, ok in tests:
        print(("PASS " if ok else "FAIL ") + n)
    print("\n%d passed, %d failed" % (len(tests) - len(fails), len(fails)))
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())

# END OF FILE
