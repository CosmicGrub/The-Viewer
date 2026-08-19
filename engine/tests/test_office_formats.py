#!/usr/bin/env python3
"""Unit tests for engine/office.py -- .docx/.xlsx/.pptx/.rtf text extraction, added while picking up
the deferred "Office format support" item (tier-gated to sysprobe.py's modern_os signal per explicit
user direction, to avoid incompatibilities on Vista/older Windows -- python-docx 1.x alone needs
Python 3.9+, ruling it out even on this app's documented Win7/Python-3.8 floor).

Covers: real extraction against real python-docx/openpyxl/python-pptx-generated files (paragraph +
table text, one-page-per-sheet, one-page-per-slide), the RTF stripper's destination-group skipping
(font/color/style tables must never leak into extracted text) against a hand-built real-shape RTF
document, the modern/legacy tier gate, and graceful degradation when a library is absent.

Format-specific sections SKIP cleanly (not fail) if their library isn't installed -- matches the
established convention for every optional dependency in this test suite (PyMuPDF/PIL/etc.).
Pure stdlib test runner; imports python-docx/openpyxl/python-pptx lazily per section."""
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import office as O


def run():
    passed, failed = [], []
    def check(name, cond):
        (passed if cond else failed).append(name)

    d = tempfile.mkdtemp(prefix="office_test_")

    # ---- modern/legacy tier gate ----------------------------------------------------------------
    orig_tier = O._modern_tier
    try:
        O._modern_tier = lambda: True
        check("modern tier: docx_available() follows the library-presence flag",
              O.docx_available() == O._DOCX_LIB_OK)
        O._modern_tier = lambda: False
        check("legacy tier: docx_available() is False regardless of library presence", O.docx_available() is False)
        check("legacy tier: xlsx_available() is False regardless of library presence", O.xlsx_available() is False)
        check("legacy tier: pptx_available() is False regardless of library presence", O.pptx_available() is False)
        check("legacy tier: extract_docx() returns '' (no crash) even with a real path",
              O.extract_docx(__file__) == "")
        check("legacy tier: extract_xlsx() returns [] (no crash)", O.extract_xlsx(__file__) == [])
        check("legacy tier: extract_pptx() returns [] (no crash)", O.extract_pptx(__file__) == [])
        check("RTF is NEVER tier-gated (dependency-free) -- rtf_available() is always True",
              O.rtf_available() is True)
    finally:
        O._modern_tier = orig_tier

    # ---- .docx --------------------------------------------------------------------------------
    try:
        import docx
        have_docx = True
    except Exception:
        have_docx = False
    if not have_docx:
        print("SKIP .docx checks (python-docx not installed)")
    else:
        doc = docx.Document()
        doc.add_paragraph("Torque the alternator bracket bolt to 45 ft-lb.")
        doc.add_paragraph("NSN 5305-01-123-4567 SCREW,MACHINE")
        tbl = doc.add_table(rows=2, cols=2)
        tbl.rows[0].cells[0].text = "Item"; tbl.rows[0].cells[1].text = "Qty"
        tbl.rows[1].cells[0].text = "Gasket"; tbl.rows[1].cells[1].text = "4"
        p = os.path.join(d, "notes.docx")
        doc.save(p)
        check("docx_available() is True with the library installed + modern tier", O.docx_available() is True)
        text = O.extract_docx(p)
        check("extract_docx(): paragraph text extracted", "Torque the alternator bracket bolt to 45 ft-lb." in text)
        check("extract_docx(): a second paragraph extracted, order preserved",
              text.index("NSN 5305-01-123-4567") > text.index("Torque the alternator"))
        check("extract_docx(): table cell text extracted too", "Gasket" in text and "Qty" in text)
        check("extract_docx(): a missing file returns '' (no crash)",
              O.extract_docx(os.path.join(d, "nope.docx")) == "")
        check("extract_docx(): None path returns '' (no crash)", O.extract_docx(None) == "")

    # ---- .xlsx --------------------------------------------------------------------------------
    try:
        import openpyxl
        have_xlsx = True
    except Exception:
        have_xlsx = False
    if not have_xlsx:
        print("SKIP .xlsx checks (openpyxl not installed)")
    else:
        wb = openpyxl.Workbook()
        ws1 = wb.active; ws1.title = "Specs"
        ws1["A1"] = "ITEM"; ws1["B1"] = "DIMENSION"
        ws1["A2"] = "Overall length"; ws1["B2"] = "180 in"
        ws2 = wb.create_sheet("Torque")
        ws2["A1"] = "Bolt"; ws2["B1"] = "45 ft-lb"
        p = os.path.join(d, "specs.xlsx")
        wb.save(p)
        check("xlsx_available() is True with the library installed + modern tier", O.xlsx_available() is True)
        sheets = O.extract_xlsx(p)
        check("extract_xlsx(): one entry per sheet, in workbook order", [s[0] for s in sheets] == ["Specs", "Torque"])
        check("extract_xlsx(): first sheet's cell text extracted", "180 in" in sheets[0][1] and "Overall length" in sheets[0][1])
        check("extract_xlsx(): second sheet is a SEPARATE page (own text, doesn't bleed into sheet 1)",
              "45 ft-lb" in sheets[1][1] and "45 ft-lb" not in sheets[0][1])
        check("extract_xlsx(): a missing file returns [] (no crash)", O.extract_xlsx(os.path.join(d, "nope.xlsx")) == [])

    # ---- .pptx --------------------------------------------------------------------------------
    try:
        import pptx
        have_pptx = True
    except Exception:
        have_pptx = False
    if not have_pptx:
        print("SKIP .pptx checks (python-pptx not installed)")
    else:
        pr = pptx.Presentation()
        s1 = pr.slides.add_slide(pr.slide_layouts[1])
        s1.shapes.title.text = "Bolt Torque"
        s1.placeholders[1].text = "Torque all bolts to 45 ft-lb before reassembly."
        s2 = pr.slides.add_slide(pr.slide_layouts[1])
        s2.shapes.title.text = "Parts List"
        s2.placeholders[1].text = "NSN 5305-01-123-4567 SCREW,MACHINE"
        p = os.path.join(d, "deck.pptx")
        pr.save(p)
        check("pptx_available() is True with the library installed + modern tier", O.pptx_available() is True)
        slides = O.extract_pptx(p)
        check("extract_pptx(): one entry per slide, 1-based, in order", [n for n, _t in slides] == [1, 2])
        check("extract_pptx(): first slide's text extracted", "Torque all bolts to 45 ft-lb" in slides[0][1])
        check("extract_pptx(): second slide is a SEPARATE page (own text, doesn't bleed into slide 1)",
              "NSN 5305-01-123-4567" in slides[1][1] and "NSN 5305-01-123-4567" not in slides[0][1])
        check("extract_pptx(): a missing file returns [] (no crash)", O.extract_pptx(os.path.join(d, "nope.pptx")) == [])

    # ---- .rtf (dependency-free, always runs) -----------------------------------------------------
    sample = (
        r"{\rtf1\ansi\deff0"
        r"{\fonttbl{\f0 Times New Roman;}{\f1 Arial;}}"
        r"{\colortbl;\red0\green0\blue0;\red255\green0\blue0;}"
        r"{\*\generator Msftedit 5.41.15.1515;}"
        r"{\stylesheet{\s0 Normal;}}"
        r"\f0\fs24 "
        r"Bolt torque 45 ft-lb required.\par "
        r"NSN 5305-01-123-4567 SCREW,MACHINE\par "
        r"}"
    )
    rtf_path = os.path.join(d, "sample.rtf")
    with open(rtf_path, "w", encoding="latin-1") as f:
        f.write(sample)
    got = O.extract_rtf(rtf_path)
    check("extract_rtf(): body text (paragraph 1) survives", "Bolt torque 45 ft-lb required." in got)
    check("extract_rtf(): body text (paragraph 2) survives", "NSN 5305-01-123-4567 SCREW,MACHINE" in got)
    check("extract_rtf(): font table names do NOT leak into the text", "Times New Roman" not in got and "Arial" not in got)
    check("extract_rtf(): the \\*\\generator ignorable destination does NOT leak (the double-push bug)",
          "Msftedit" not in got)
    check("extract_rtf(): stylesheet name does NOT leak", "Normal" not in got)
    check("extract_rtf(): a missing file returns '' (no crash)", O.extract_rtf(os.path.join(d, "nope.rtf")) == "")
    # a genuinely present but malformed/unbalanced file must not raise
    bad_path = os.path.join(d, "bad.rtf")
    with open(bad_path, "wb") as f:
        f.write(b"not rtf at all, just plain bytes {{{ unbalanced")
    try:
        r = O.extract_rtf(bad_path)
        ok = isinstance(r, str)
    except Exception:
        ok = False
    check("extract_rtf(): malformed/unbalanced RTF never raises", ok)

    return passed, failed


if __name__ == "__main__":
    p, f = run()
    for n in p: print("PASS", n)
    for n in f: print("FAIL", n)
    print("\n%d passed, %d failed" % (len(p), len(f)))
    sys.exit(1 if f else 0)

# END OF FILE
