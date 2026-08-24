#!/usr/bin/env python3
"""THE VIEWER -- Office document text extraction (.docx/.xlsx/.pptx/.rtf), tier-gated.

python-docx/openpyxl/python-pptx are real, not-stdlib dependencies (requirements.txt) -- and this
app's design bar is COMPLETE feature parity down to Windows Vista/7 with no accuracy loss (see
docs/SYSTEM-REQUIREMENTS.md), the same reasoning that already keeps GPU-accelerated OCR a Win10+
-only *speed* extra rather than a requirement. Office parsing follows the identical pattern: gated
on sysprobe.py's own modern_os signal (Win10/11, osrank>=100 -- the SAME tier test the render/OCR
engine substitution and camelot_tables()'s backend selection already use), not just library
presence, so a legacy machine that happens to have one of these pip packages installed still
doesn't pay for/attempt something this app's design bar doesn't ask of it. On the legacy tier (or
if a library genuinely isn't installed, even on the modern tier) these formats are still DISCOVERED
by classify_ext()/crawl() exactly as before this module existed -- just 0 extracted pages, the same
"office" fallback behavior that predates the Discovery Engine entirely. Each format degrades
independently (docx_available()/xlsx_available()/pptx_available()) -- one broken/missing library
must never take the other two (or the dependency-free RTF path) down with it.

.doc/.ppt/.xls (pre-2007 binary Office formats) remain unsupported everywhere -- python-docx/
openpyxl/python-pptx only read the modern XML-based formats; the legacy binary formats would need
yet more dependencies (xlrd for .xls; there is no good pure-Python .doc/.ppt reader at all) for
formats a military-TM corpus is unlikely to hold much of in the first place. Discovered, 0 pages,
same as today.

RTF (.rtf) is the one format here that's genuinely dependency-free -- no good stdlib-free-and-
lightweight RTF library exists, so this ships its own small, tolerant stripper (same "good enough
for a structured export, not meant to survive adversarial/malformed markup" bar as
viewer_ingest.py's _strip_html()), no tier gate needed since it never imports anything new."""
import os
import re

try:
    import docx as _docx
    _DOCX_LIB_OK = True
except Exception:
    _docx = None
    _DOCX_LIB_OK = False

try:
    import openpyxl as _openpyxl
    _XLSX_LIB_OK = True
except Exception:
    _openpyxl = None
    _XLSX_LIB_OK = False

try:
    import pptx as _pptx
    _PPTX_LIB_OK = True
except Exception:
    _pptx = None
    _PPTX_LIB_OK = False


def _modern_tier():
    """Same modern/legacy signal every other OS-substituted engine in this codebase uses
    (sysprobe.py's build_profile()['modern_os']). Fails OPEN (assume modern, attempt the real
    import) if the probe itself can't run for any reason -- same precedent tables_plus.py's
    _camelot_backend() already sets ("a probe glitch must never break extraction"). This is safe
    specifically because the try/except import guards above are the REAL protection against a
    genuinely incompatible environment (an old Python failing to even parse a modern library's
    syntax raises a plain ImportError/SyntaxError there, caught already) -- this tier gate's job is
    purely "don't bother trying on a machine we already know is legacy", not crash prevention."""
    try:
        here = os.path.dirname(os.path.abspath(__file__))
        import sys as _sys
        if here not in _sys.path:
            _sys.path.insert(0, here)
        import sysprobe
        return bool(sysprobe.load_or_build().get("modern_os", True))
    except Exception:
        return True


def docx_available():
    return _DOCX_LIB_OK and _modern_tier()


def xlsx_available():
    return _XLSX_LIB_OK and _modern_tier()


def pptx_available():
    return _PPTX_LIB_OK and _modern_tier()


def rtf_available():
    return True   # dependency-free -- no tier gate needed, see module docstring


def extract_docx(path):
    """Whole-document text as ONE string -- Word's XML doesn't record page boundaries (the same
    "no native page concept" reasoning viewer_ingest.py's index_other() already applies to .txt/
    .html, which also become a single page there). Includes table cell text (a real TM-style parts/
    spec table embedded in a Word doc is common), paragraph order preserved. Returns "" on any
    failure -- never raises, matches every other extractor in this codebase's contract."""
    if not docx_available() or not path or not os.path.exists(path):
        return ""
    try:
        d = _docx.Document(path)
        parts = [p.text.strip() for p in d.paragraphs if p.text and p.text.strip()]
        for t in d.tables:
            for row in t.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        return ""


def extract_xlsx(path):
    """[(sheet_name, text), ...] -- ONE entry per sheet, a natural page-per-sheet mapping (a multi-
    sheet workbook becomes a multi-page document, mirroring how a multi-page PDF already works,
    unlike .docx's single-page treatment above). Cell values tab-joined per row, rows newline-
    joined. read_only=True keeps memory bounded on a large workbook. Returns [] on any failure."""
    if not xlsx_available() or not path or not os.path.exists(path):
        return []
    out = []
    try:
        wb = _openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            for ws in wb.worksheets:
                lines = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                    if cells:
                        lines.append("\t".join(cells))
                out.append((ws.title or "Sheet", "\n".join(lines)))
        finally:
            wb.close()
    except Exception:
        return []
    return out


def extract_pptx(path):
    """[(slide_number, text), ...], 1-based -- pptx has explicit slide boundaries, a genuinely
    natural page mapping (unlike .docx). Every text frame + table cell on the slide, in shape
    order. Returns [] on any failure."""
    if not pptx_available() or not path or not os.path.exists(path):
        return []
    out = []
    try:
        pr = _pptx.Presentation(path)
        for i, slide in enumerate(pr.slides, start=1):
            parts = []
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text.strip())
                    elif getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [(c.text or "").strip() for c in row.cells]
                            cells = [c for c in cells if c]
                            if cells:
                                parts.append(" | ".join(cells))
                except Exception:
                    continue
            out.append((i, "\n".join(parts)))
    except Exception:
        return []
    return out


# ---- RTF: dependency-free, tolerant stripper -----------------------------------------------------
# A hand-rolled brace-depth walk, not pure regex -- RTF's nesting isn't a regular language, and a
# naive "strip every \controlword" regex would leak font-table/color-table/stylesheet NAMES straight
# into the extracted text (every real-world RTF file has at least a font table; this is the norm,
# not an edge case) -- a meaningful quality loss for both search relevance and dimensional-extraction
# false positives (a color table's numeric RGB triples read exactly like measurement noise).
_RTF_CTRLWORD_RE = re.compile(r"([a-zA-Z]+)(-?\d+)?[ ]?")
_RTF_HEX_RE = re.compile(r"'([0-9a-fA-F]{2})")
# Known destination groups to always drop entirely, even without a leading \* marker (per the RTF
# spec, these are recognized by name alone) -- \* IS handled generically below for anything else.
_RTF_DESTINATIONS = {"fonttbl", "colortbl", "stylesheet", "pict", "object", "footnote",
                      "header", "headerf", "headerl", "headerr", "footer", "footerf", "footerl",
                      "footerr", "info", "generator", "listtable", "listoverridetable", "revtbl",
                      "themedata", "colorschememapping", "datastore", "xmlnstbl", "rsidtbl"}


def extract_rtf(path, max_bytes=5_000_000):
    """Dependency-free RTF -> plain text. Strips control words/groups, skips ignorable destination
    groups entirely (font/color/style tables, embedded pictures/objects -- both the RTF-spec \\*
    generic marker and the common named destinations above), decodes \\'xx hex escapes as latin-1
    bytes (RTF's default codepage for plain ASCII/Western text), and turns \\par/\\line into
    newlines. Not a real RTF parser -- good enough for a structured document export, not meant to
    survive adversarial/malformed RTF. Returns "" on any failure."""
    if not path or not os.path.exists(path):
        return ""
    try:
        with open(path, "rb") as f:
            raw = f.read(max_bytes)
        text = raw.decode("latin-1", errors="ignore")
    except Exception:
        return ""
    try:
        out = []
        i, n, depth = 0, len(text), 0
        skip_depths = []   # stack of depths currently inside an ignorable destination group
        while i < n:
            ch = text[i]
            if ch == "{":
                depth += 1
                i += 1
                continue
            if ch == "}":
                if skip_depths and skip_depths[-1] == depth:
                    skip_depths.pop()
                depth -= 1
                i += 1
                continue
            if ch == "\\":
                if text[i:i + 2] == "\\*":
                    # generic ignorable-destination marker (RTF spec): whatever named control word
                    # follows (recognized or not), the group it's in should be skipped entirely.
                    skip_depths.append(depth)
                    i += 2
                    continue
                if text[i:i + 2] in ("\\{", "\\}", "\\\\"):
                    if not skip_depths:
                        out.append(text[i + 1])
                    i += 2
                    continue
                mh = _RTF_HEX_RE.match(text, i + 1)
                if mh:
                    if not skip_depths:
                        out.append(chr(int(mh.group(1), 16)))
                    i = mh.end()
                    continue
                mc = _RTF_CTRLWORD_RE.match(text, i + 1)
                if mc:
                    word = mc.group(1)
                    # only push a NEW skip level for a NAMED destination if we didn't just push one
                    # via an immediately-preceding \* marker for this exact group -- \*\generator
                    # (a \* marker followed by a name that ALSO happens to be in _RTF_DESTINATIONS)
                    # would otherwise double-push for the same group, leaving a permanently stale
                    # entry that never gets popped once its real closing brace only removes one.
                    already_marked = bool(skip_depths) and skip_depths[-1] == depth
                    if word in _RTF_DESTINATIONS and not already_marked:
                        skip_depths.append(depth)
                    elif not skip_depths and word in ("par", "line", "row"):
                        out.append("\n")
                    elif not skip_depths and word == "tab":
                        out.append("\t")
                    i = mc.end()
                    continue
                i += 1   # a lone/unrecognized backslash -- drop it, never emit
                continue
            if not skip_depths:
                out.append(ch)
            i += 1
        result = "".join(out)
        result = re.sub(r"[ \t]+", " ", result)
        result = re.sub(r"\n[ \t]*\n+", "\n\n", result)
        return result.strip()
    except Exception:
        return ""


if __name__ == "__main__":
    # Self-test: a hand-built, real-shape minimal RTF document (font table + color table + a
    # generic \* ignorable destination + body text with a \par break) -- proves destination groups
    # never leak into the extracted text and real body content survives intact.
    sample = (
        r"{\rtf1\ansi\deff0"
        r"{\fonttbl{\f0 Times New Roman;}{\f1 Arial;}}"
        r"{\colortbl;\red0\green0\blue0;\red255\green0\blue0;}"
        r"{\*\generator Msftedit 5.41.15.1515;}"
        r"{\stylesheet{\s0 Normal;}}"
        r"\f0\fs24 "
        r"Bolt torque 45 ft-lb required.\par "
        r"NSN 5305-01-123-4567 SCREW,MACHINE\par "
        r"}"
    )
    import tempfile
    p = os.path.join(tempfile.mkdtemp(), "sample.rtf")
    with open(p, "w", encoding="latin-1") as f:
        f.write(sample)
    got = extract_rtf(p)
    assert "Bolt torque 45 ft-lb required." in got, ("body text missing", got)
    assert "NSN 5305-01-123-4567 SCREW,MACHINE" in got, ("second paragraph missing", got)
    assert "Times New Roman" not in got and "Arial" not in got, ("font table leaked", got)
    assert "Msftedit" not in got, ("generator destination leaked", got)
    assert "Normal" not in got, ("stylesheet leaked", got)
    print("office.py RTF self-test OK -- destinations skipped, body text intact:")
    print(repr(got))

# END OF FILE
